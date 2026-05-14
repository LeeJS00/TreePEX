"""make_ppt.py — generate presentation .pptx for TreePEX work.

Builds 14 slides covering: motivation, baseline, diagnosis, hypothesis,
H1-H5 experiments, Tweedie loss, auto-4pct, deployment demo, comparison
to prior, conclusions.

Output: TreePEX/presentation/TreePEX_presentation.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path("/home/jslee/projects/PINNPEX/TreePEX/presentation")
FIGS = ROOT / "figures"

PRES_W = 13.333  # 16:9 widescreen
PRES_H = 7.5

prs = Presentation()
prs.slide_width = Inches(PRES_W)
prs.slide_height = Inches(PRES_H)


# --- helpers ----------------------------------------------------------

def add_blank_slide():
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def add_title(slide, text: str, *, size: int = 32, color=(13, 71, 161)):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(PRES_W - 1), Inches(0.7))
    tf = box.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size); p.font.bold = True
    p.font.color.rgb = RGBColor(*color)


def add_subtitle(slide, text: str, *, size: int = 16, top: float = 0.95):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(PRES_W - 1), Inches(0.5))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size); p.font.italic = True
    p.font.color.rgb = RGBColor(80, 80, 80)


def add_bullets(slide, bullets, *, left=0.6, top=1.6, width=12.0, height=4.5,
                font=14, line_spacing=1.2):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p.text = ("• " if level == 0 else "  – ") + text
        p.level = 0
        p.font.size = Pt(font - level)
        p.line_spacing = line_spacing


def add_image(slide, fig_name: str, *, left, top, width=None, height=None):
    p = FIGS / fig_name
    if not p.exists():
        return
    if width and height:
        slide.shapes.add_picture(str(p), Inches(left), Inches(top),
                                  width=Inches(width), height=Inches(height))
    elif width:
        slide.shapes.add_picture(str(p), Inches(left), Inches(top), width=Inches(width))
    elif height:
        slide.shapes.add_picture(str(p), Inches(left), Inches(top), height=Inches(height))
    else:
        slide.shapes.add_picture(str(p), Inches(left), Inches(top))


def add_table(slide, data, *, left=0.6, top=2.0, width=12.0, height=2.8,
              header_color=(13, 71, 161), highlight_rows=None):
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                    Inches(width), Inches(height)).table
    highlight_rows = highlight_rows or []
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)
                    run.font.name = "Calibri"
                if r == 0:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.bold = True; run.font.color.rgb = RGBColor(255, 255, 255)
                if r in highlight_rows:
                    for run in paragraph.runs:
                        run.font.bold = True
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*header_color)
            elif r in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(220, 240, 220)


def add_footer(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(PRES_H - 0.4),
                                    Inches(PRES_W - 1), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9); p.font.italic = True
    p.font.color.rgb = RGBColor(120, 120, 120)


# === SLIDES ===========================================================

# --- 1. Title ---
s = add_blank_slide()
box = s.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(PRES_W - 2), Inches(1.4))
p = box.text_frame.paragraphs[0]
p.text = "Practical Cross-Design Cap Prediction"
p.font.size = Pt(40); p.font.bold = True
p.font.color.rgb = RGBColor(13, 71, 161); p.alignment = PP_ALIGN.CENTER
p2 = box.text_frame.add_paragraph()
p2.text = "120× Faster than PINN at 4.98 % MAPE on tv80s"
p2.font.size = Pt(26); p2.font.italic = True
p2.font.color.rgb = RGBColor(60, 60, 60); p2.alignment = PP_ALIGN.CENTER

box = s.shapes.add_textbox(Inches(1), Inches(4.8), Inches(PRES_W - 2), Inches(2))
tf = box.text_frame
tf.paragraphs[0].text = "Top-K aggressor features + Tweedie loss + 5-seed ensemble"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = ""
p2 = tf.add_paragraph()
p2.text = "PINN-PEX project · TreePEX deployment demo · 2026-05-10"
p2.font.size = Pt(14); p2.font.italic = True
p2.font.color.rgb = RGBColor(120, 120, 120); p2.alignment = PP_ALIGN.CENTER


# --- 2. Problem statement ---
s = add_blank_slide()
add_title(s, "1. Problem: ML-based parasitic extraction (PEX)")
add_subtitle(s, "Predict per-net capacitance from routed layout — fast enough for production")
add_bullets(s, [
    "VLSI design flow needs accurate net-level cap (gnd / cpl) for IR-drop, timing, crosstalk analysis",
    ("Golden oracle: Synopsys StarRC — minutes per chip, license-bound", 1),
    ("Pattern-matching tools (Innovus, OpenRCX): seconds per chip but 22-72 % MAPE per cap-decile", 1),
    ("ML approach goal: <10 % MAPE in <1 s, license-free, reproducible", 1),
    "",
    "Prior ML-PEX work (ParaGraph, GNNCap, our v12 PINN):",
    ("v12 PINN (5-seed): tv80s tot 5.55 %, gnd 22.59 %, cpl 17.53 % at 20.4 s standalone", 1),
    ("Per-channel MAPE 17–22 % is too high for downstream IR-drop sensitivity", 1),
    "",
    "Question: can we beat v12 PINN on accuracy AND speed simultaneously?",
])


# --- 3. v12 PINN baseline ---
s = add_blank_slide()
add_title(s, "2. Baseline: v12 PINN (44K params, hybrid mesh)")
add_subtitle(s, "Cuboid-set encoder + neural flux router — what we'll compare against")
add_bullets(s, [
    "Architecture: per-net cuboid set → CuboidEncoder (DeepSet) → NeuralFluxRouter → 3 heads (gnd, cpl, screening)",
    "Training: Phase 1 SSL pretrain → Phase 2 active-learning fine-tune against StarRC golden",
    "Joint Pareto v12 5-seed (locked frontier as of 2026-05-04):",
    ("tv80s tot_med 5.55 %, gnd 22.59 %, cpl 17.53 %", 1),
    ("Per-net wall-clock: 20.4 s standalone on tv80s 3,169 nets (~6 ms/net)", 1),
    "",
    "Limitations identified from v12 residuals:",
    ("Per-channel gap: gnd 22.6 % is 4× larger than B1 XGBoost's gnd-channel error", 1),
    ("DeepSet aggregation symmetrizes cuboid features — loses per-aggressor identity", 1),
    ("Cuboid encoder operates on TARGET cuboids only (no aggressor features)", 1),
])


# --- 4. Diagnosis: within-bucket R² ---
s = add_blank_slide()
add_title(s, "3. Diagnosis (H1, H2): within-bucket R² is the bottleneck")
add_subtitle(s, "Overall R² 0.99 hides the real issue")
add_bullets(s, [
    ("Overall R² across all nets: 0.9954 — looks excellent…", 0),
    ("But stratify by c_total cap-decile: within-bucket R² is poor in mid range", 0),
    ("3-way oracle blend (v12+B1+PA-XGB) gain: only +0.01–0.03 R² in mid bucket", 0),
    "",
    ("Big XGBoost (depth=12, n_est=3000) ALONE: nova MAPE regress (5.70 → 6.30)", 0),
    ("Conclusion: this is an INFORMATION ceiling on DEF/LEF/Liberty inputs, not capacity", 0),
    "",
    "Hypothesis: missing per-aggressor identity (current 41 features have only aggregates)",
], top=1.4, height=2.5)
add_image(s, "fig1_per_bucket_R2.png", left=0.5, top=4.0, width=12.3)
add_footer(s, "fig1: per-bucket cpl R² across cap-deciles. Top-K features (orange/green) lift mid-bucket discrimination ~+0.1 R²")


# --- 5. Hypothesis: Top-K aggressor features ---
s = add_blank_slide()
add_title(s, "4. Hypothesis: Top-K aggressor features (H3)")
add_subtitle(s, "Per-aggressor identity → per-bucket discrimination")
add_bullets(s, [
    "Current 41 features (NetFeatureVector) capture only AGGREGATE aggressor stats:",
    ("e.g., n_aggressor_nets, broadside_overlap_total, broadside_overlap_p95", 1),
    ("Top-1 aggressor's individual overlap, distance, layer-diff are LOST in aggregation", 1),
    "",
    "New features (26 per net, extracted from raw tile pickles in 3.5 h on 8 workers):",
    ("Per top-K aggressor (k=1,2,3): score, overlap_um2, min_xy_dist_um, mean_dz_um, agg_size_um2, layer_diff_flag", 1),
    ("Aux: agg_count_above/below_target_z, agg_n_distinct, topk_score_concentration", 1),
    ("3D density: agg_count_within_{1,3,5}μm radius from target centroid", 1),
    "",
    "Physics-motivated score: Σᵢⱼ ε · xy_overlap / max(|Δz|, ε_z) over (target, agg) cuboid pairs",
    "",
    "Pipeline: pex_v3 cuboid store + aggressor groups → vectorized numpy pair tensors → score + sort → top-K",
])


# --- 5b. Feature pipeline diagram ---
s = add_blank_slide()
add_title(s, "5a. Methodology: feature pipeline (DEF→67-D vector)")
add_subtitle(s, "Two precomputed CSVs joined per (design, net) — 221k × 67 features")
add_image(s, "fig6_feature_pipeline.png", left=0.2, top=1.2, width=13.0, height=5.5)
add_footer(s, "fig6: pipeline architecture. Left: raw inputs. Middle: parsers + extractors. Right: joined feature table fed to XGBoost.")

# --- 5c. Feature catalog detail ---
s = add_blank_slide()
add_title(s, "5b. Feature catalog (67 dims total)")
add_subtitle(s, "41 base + 26 H3 top-K aggressor — distinct semantic groups")
add_bullets(s, [
    "Base 41 (NetFeatureVector — frozen schema):",
    ("Geometric (6): total_wire_length_um, total_metal_area_um2, n_cuboids, bbox_xy_um2, bbox_z_um, aspect_ratio", 1),
    ("Layer histogram (9): layer_hist_M1 ... M9_plus", 1),
    ("Coupling aggregate (11): n_aggressor_nets, broadside/lateral_overlap_total/p95, spacing_min/p25/p50/p95, n_edges_lt_1um/1_to_3um/3_to_4um", 1),
    ("VSS shielding (5): vss_n_cuboids, vss_total_metal_area_um2, vss_shield_M1_M3 / M4_M5 / M6_plus", 1),
    ("Layer stack (5): eps_min/max/mean, n_layers_present, fanout", 1),
    ("Density per layer group (3): density_M1_M3 / M4_M5 / M6_plus", 1),
    ("Compact-model intermediates (2): compact_gnd_estimate_fF, compact_cpl_estimate_total_fF", 1),
    "",
    "New 26 (H3 top-K aggressor — paper contribution):",
    ("Top-K (k=1,2,3) per-aggressor (18): score, overlap_um2, min_xy_dist_um, mean_dz_um, agg_size_um2, layer_diff_flag", 1),
    ("Aux (4): agg_count_above/below_target_z, agg_n_distinct, topk_score_concentration", 1),
    ("3D density (3): agg_count_within_{1,3,5}μm xyz radius from target centroid", 1),
    ("Sanity (1): target_n_cuboids_check", 1),
    "",
    "Score formula (physics-motivated parallel-plate proxy):",
    ("    score(agg) = Σ over (target, agg) cuboid pairs of [ε_avg · xy_overlap / max(|Δz|, 0.05 μm)]", 1),
    ("Aggressors ranked by score; top-3 carry the discrimination signal.", 1),
], font=11, line_spacing=1.05, height=5.5)

# --- 5d. Model architecture ---
s = add_blank_slide()
add_title(s, "5c. Model architecture: 5-seed Tweedie XGBoost ensemble")
add_subtitle(s, "Two heads (gnd / cpl) × 5 seeds = 10 weight files (~12 MB each)")
add_image(s, "fig7_xgboost_architecture.png", left=0.2, top=1.2, width=13.0, height=5.5)
add_footer(s, "fig7: (left) per-seed boosting structure (depth 8, 500 trees max with early-stop). (right) Tweedie objective derivation + comparison vs alternatives.")

# --- 5e. Training methodology ---
s = add_blank_slide()
add_title(s, "5d. Training methodology + 5-seed protocol")
add_subtitle(s, "Manifest H1 split (no design-net leakage), early-stopping on valid log-MSE")
add_image(s, "fig8_training_protocol.png", left=0.2, top=1.2, width=13.0, height=5.5)
add_footer(s, "fig8: (top) 51/5.7/43 train/valid/test split. (bottom) 5-seed protocol + hyperparameters + Tweedie objective.")

# --- 5f. Inference details ---
s = add_blank_slide()
add_title(s, "5e. Inference: ensemble averaging at PREDICTION level")
add_subtitle(s, "predict-then-aggregate cancels per-seed variance — 0.108 pp gain over 5-seed-mean-of-MAPE")
add_bullets(s, [
    "Two ways to combine 5 trained seeds:",
    "",
    "(A) Average MAPE across seeds (= traditional 5-seed lock):",
    ("seed_42_MAPE = 5.05, seed_0 = 5.13, seed_1 = 5.09, seed_2 = 5.10, seed_3 = 5.07", 1),
    ("→ mean = 5.087 ± 0.049  ← S4 Tweedie 5-seed lock value", 1),
    "",
    "(B) Average PREDICTIONS across seeds, then compute MAPE (= TreePEX ensemble):",
    ("pred_gnd_ens(net) = mean over 5 seeds of pred_gnd_seed(net)", 1),
    ("pred_cpl_ens(net) = mean over 5 seeds of pred_cpl_seed(net)", 1),
    ("MAPE_ens(test) = MAPE(pred_ens, gold)", 1),
    ("→ tv80s tot_med = 4.979, nova = 5.279", 1),
    "",
    "Why (B) is better:",
    ("Per-seed model has independent stochastic variance from subsample 0.8 + colsample 0.8", 1),
    ("Averaging predictions FIRST cancels noise BEFORE it hits |p−y|/y nonlinearity", 1),
    ("Effectively a Bayesian model average — same cost as a single model at inference", 1),
    "",
    "Inference computation per net:",
    ("  pred_gnd = (1/5) · Σ_seed XGB_gnd_seed.predict(features)  ← 5 tree traversals", 1),
    ("  pred_cpl = (1/5) · Σ_seed XGB_cpl_seed.predict(features)", 1),
    ("  Total: 10 XGB_predict calls per net, ~5 μs/net on CPU (vectorized over batch)", 1),
], font=11, line_spacing=1.1, height=5.7)

# --- 6. H4 + H5 experiments ---
s = add_blank_slide()
add_title(s, "6. Validation H4-H5: features × capacity synergy")
add_subtitle(s, "Both must scale together — capacity alone overfits")
data = [
    ["config", "feats", "tv80s tot_med", "tv80s gnd", "tv80s cpl", "nova tot_med", "inference"],
    ["B1 (5-seed) — base", "41", "5.30 ± 0.05", "19.89 %", "14.16 %", "5.83 ± 0.08", "0.05 s"],
    ["Small + new features", "67", "5.28 ± 0.07", "18.90 %", "13.57 %", "5.62 ± 0.02", "0.05 s"],
    ["Big + 41 features only", "41", "5.31 ± 0.05", "18.95 %", "13.99 %", "6.24 ± 0.07 ⚠", "0.4 s"],
    ["Big + new features", "67", "5.17 ± 0.07", "17.84 %", "13.23 %", "5.92 ± 0.07 ⚠", "0.4 s"],
]
add_table(s, data, left=0.5, top=1.5, width=12.3, height=2.4, highlight_rows=[2])
add_bullets(s, [
    "H2 (Big alone, 41 feat): nova regression 5.70→6.24 confirmed structural — capacity OVERFITS without new features",
    "H4 (Small + new): per-channel gnd −1.0 pp, cpl −0.6 pp; tot stable; Mid-bucket C2-C8 cpl R² +0.07 vs B1",
    "H5 (Big + new): synergy — capacity exploits new feature dimensions; tv80s tot −0.13 pp; mid-bucket cpl R² +0.10",
    "Pattern: capacity helps when feature space expands meaningfully (H2 isolated this)",
], top=4.2, height=2.8, font=12)


# --- 7. Tweedie loss ---
s = add_blank_slide()
add_title(s, "7. Loss alignment (S4): Tweedie variance_power=1.5")
add_subtitle(s, "MAPE-aligned objective for power-law cap distribution")
add_bullets(s, [
    "Default objective `reg:squarederror` minimizes MSE on log1p target — NOT MAPE-aligned",
    "Tweedie (variance_power=1.5): power-law cap distribution + log link — closer to MAPE",
    ("XGBoost: `objective='reg:tweedie'`, `tweedie_variance_power=1.5`", 1),
    "",
    "S4 Tweedie 5-seed (Small_combined config: depth=8, n_est=500, lr=0.05):",
    ("tv80s: 5.087 ± 0.049 / gnd 18.02 / cpl 13.35  (vs Small_combined squarederror 5.28)", 1),
    ("nova:  5.417 ± 0.027 / gnd 17.48 / cpl 15.00  (vs squarederror 5.62, std 4× lower)", 1),
    "",
    "vp grid (S9): 1.2 / 1.4 / 1.6 / 1.8 — best vp = 1.4 single-seed (5.032), but variance pushes 5-seed mean above 5.087",
    "Big + Tweedie (S6): WORSE than Small + Tweedie — capacity hurts even with aligned loss",
    "",
    "Quantile loss (P1, α=0.5): 5.380 / 6.704 — worse; |p-y| anchor ≠ MAPE anchor",
    "Custom MAPE objective (P3): unstable, immediate early-stop with constant-Hessian sign-gradient",
])


# --- 8. Frontier evolution ---
s = add_blank_slide()
add_title(s, "8. Frontier evolution: v12 PINN → TreePEX ensemble")
add_subtitle(s, "Step-by-step accuracy + runtime improvement")
add_image(s, "fig2_5seed_progression.png", left=0.2, top=1.4, width=13.0)
add_footer(s, "fig2: (a) 5-seed mean tv80s + nova MAPE per method (b) inference wall on tv80s 3,169 nets (log scale)")


# --- 9. Auto-4pct exploration ---
s = add_blank_slide()
add_title(s, "9. Auto-4pct: 14 strategies × goal of tv80s ≤ 4 %")
add_subtitle(s, "Exhaustive ladder over loss / blends / per-bucket / recursive — best 4.979 (deployable)")
add_image(s, "fig3_strategy_ladder.png", left=0.5, top=1.4, width=12.3, height=5.5)
add_footer(s, "fig3: 14-strategy comparison. Green = deployable, gray = oracle blends, red = NOT-deployable oracle (P2 / P7)")


# --- 10. Per-bucket TreePEX result ---
s = add_blank_slide()
add_title(s, "10. TreePEX per-bucket: C1 noise floor, C8 = 4.02 %")
add_subtitle(s, "Aggregate 4.98 % is bottlenecked by C1-C2 (cap < 0.20 fF) denominator noise")
add_image(s, "fig4_per_bucket_TreePEX.png", left=0.5, top=1.4, width=12.3, height=5.0)
add_footer(s, "fig4: TreePEX ensemble per-cap-decile MAPE on tv80s (blue) and nova (orange). Mid-bucket C8 (cap mean 1.46 fF) hits 4.02 % on tv80s")


# --- 11. TreePEX deployment demo ---
s = add_blank_slide()
add_title(s, "11. TreePEX deployment: end-to-end SPEF tool")
add_subtitle(s, "fresh-clone reproducible — 10 weight files (120 MB total) + 3 stage scripts")
add_bullets(s, [
    "Stage 1 [02_inference.py]: features → 5-seed ensemble predict (mean) → CSV — 0.171 s tv80s",
    "Stage 2 [03_write_spef.py]: predictions → IEEE 1481-1999 SPEF — 0.16 s tv80s, 3.0 s nova",
    "Stage 3 [04_compare_golden.py]: parse pred + golden SPEF, align by net, compute MAPE / R²",
    "Orchestrator [pex_tool.py]: --design X | --all (both designs)",
], top=1.3, height=2.0, font=13)
data = [
    ["design", "n_nets", "tot_med", "gnd_med", "cpl_med", "R²(tot)", "inference", "SPEF write"],
    ["tv80s", "3,169",   "4.979 %", "18.02 %", "13.27 %", "0.9940", "0.171 s", "0.16 s"],
    ["nova",  "92,425",  "5.279 %", "17.40 %", "14.96 %", "0.9911", "0.185 s", "3.01 s"],
]
add_table(s, data, left=0.5, top=3.5, width=12.3, height=1.3, highlight_rows=[1, 2])
add_bullets(s, [
    "SPEF round-trip: max abs err 5e-6 fF (lossless)",
    "Ensemble inference at ~5 μs/net throughput (large-design batching amortizes overhead)",
    "vs v12 PINN: tv80s tot −0.57 pp at 120× faster wall (0.17 s vs 20.4 s)",
], top=5.2, height=1.6, font=13)


# --- 12. Speed-accuracy Pareto ---
s = add_blank_slide()
add_title(s, "12. Speed-accuracy Pareto: TreePEX dominant point")
add_subtitle(s, "Faster than B1 XGBoost-comparable, lower MAPE than v12 PINN")
add_image(s, "fig5_speed_accuracy_pareto.png", left=1.5, top=1.3, width=10.0, height=5.5)
add_footer(s, "fig5: log-log Pareto. TreePEX (★) sits below v12 (■) and beside B1 (●). P2 oracle (✕) shows non-deployable upper bound.")


# --- 13. Comparison vs prior ML-PEX ---
s = add_blank_slide()
add_title(s, "13. Comparison vs prior ML-PEX papers")
add_subtitle(s, "Each axis: contribution / performance / novelty")
data = [
    ["method", "input", "model", "tv80s tot", "inference 95k", "novelty"],
    ["ParaGraph (DAC '21)",  "DEF + tech LEF",       "GNN attention",       "~5–10 %",   "~5–10 s",  "graph attention for layout"],
    ["GNNCap (DAC '22)",      "DEF + tech LEF",       "GNN",                 "~5–7 %",    "~3–5 s",   "training-time efficiency"],
    ["v12 PINN (ours, 5-seed)","DEF+LEF+Liberty",     "cuboid-enc + flux",   "5.55 %",    "20.4 s",   "physics-informed mesh"],
    ["B1 XGBoost (5-seed)",   "DEF+LEF+Liberty",      "XGB on 41 hand feat", "5.30 %",    "0.05 s",   "hand-feature baseline"],
    ["TreePEX (this work)",     "DEF+LEF+Liberty",     "5-seed Tweedie XGB ens", "4.979 %",  "0.36 s",   "top-K aggressor + Tweedie + ensemble"],
]
add_table(s, data, left=0.3, top=1.5, width=12.7, height=2.6, highlight_rows=[5])
add_bullets(s, [
    "Performance: TreePEX beats every prior on tv80s tot at competitive or faster wall",
    "Contribution: practical deployable tool with end-to-end SPEF write, not just research model",
    "Novelty: top-K aggressor features (vs aggregate-only) + Tweedie loss + within-bucket diagnostic framework",
    "Risk: simpler model than ParaGraph/GNNCap → counter with multi-pillar narrative (features + loss + ceiling)",
], top=4.4, height=2.6, font=12)


# --- 14. Conclusions + future ---
s = add_blank_slide()
add_title(s, "14. Conclusions + future work")
add_subtitle(s, "What we learned, what's still open")
add_bullets(s, [
    "Frontier: TreePEX 5-seed Tweedie XGBoost ensemble — tv80s 4.979 %, nova 5.279 %, 0.17 s inference",
    ("Beats v12 PINN by 0.57 pp at 120× faster end-to-end wall", 1),
    ("End-to-end SPEF pipeline reproducible from fresh clone (120 MB weights)", 1),
    "",
    "Three contribution pillars:",
    ("(1) Top-K aggressor features — per-aggressor identity that prior aggregate-stats methods miss", 1),
    ("(2) Tweedie loss alignment — 0.1-0.3 pp MAPE-aligned gain", 1),
    ("(3) Within-bucket R² + 4-way oracle bound — methodological diagnostic framework", 1),
    "",
    "4 % goal NOT closed by 14 strategies under 67-D scalar feature + tree formulation",
    ("Per-bucket C8 already at 4.02 % locally — aggregate bottleneck is C1 (cap < 0.15 fF), but…", 1),
    ("…StarRC golden uses the SAME inputs we do (DEF + LEF + Liberty + layer-stack)", 1),
    ("…→ The 1.0 pp gap is REPRESENTATION-side, not input-side", 1),
    "",
    "Path to 4 % (model-side; inputs unchanged):",
    ("Per-pair LEARNED coupling head (vs failed Strike #2 which used uniform analytic)", 1),
    ("Pattern-bank memory / k-NN inference (mimics StarRC NXTGRD lookup)", 1),
    ("Stack v12 encoder representation as features for XGBoost residual", 1),
    ("Different output target (block-level cap, length-weighted MAPE)", 1),
])


# --- save ---
out_path = ROOT / "TreePEX_presentation.pptx"
prs.save(out_path)
print(f"[ok] saved {out_path}  ({len(prs.slides)} slides)")
